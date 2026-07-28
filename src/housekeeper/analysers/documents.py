import hashlib
import re
import unicodedata
from collections.abc import Iterable
from pathlib import Path


def normalize_document_text(text: str, max_chars: int) -> str:
    return re.sub(r"\s+", " ", unicodedata.normalize("NFKC", text.replace("\x00", " "))).strip()[
        :max_chars
    ]


def compute_normalized_text_hash(text: str) -> str:
    return hashlib.sha256(text.encode()).hexdigest()


def extract_plaintext_metadata(path: Path, config):
    try:
        maximum = config.section("documents")["max_text_characters"]
        with path.open("r", encoding="utf-8", errors="replace") as handle:
            raw = handle.read(maximum * 4)
        text = normalize_document_text(
            raw,
            maximum,
        )
        return {
            "normalized_text": text,
            "character_count": len(text),
            "word_count": len(text.split()),
            "normalized_text_hash": compute_normalized_text_hash(text),
            "extraction_status": "OK",
        }
    except OSError as exc:
        return {"extraction_status": "ERROR", "extraction_error": str(exc)}


def _bounded_text(parts: Iterable[str], config):
    maximum = config.section("documents")["max_text_characters"]
    return normalize_document_text("\n".join(parts), maximum)


def _structured_document(path: Path, suffix: str, config):
    try:
        if suffix == "docx":
            from docx import Document  # type: ignore[import-not-found]

            document = Document(str(path))
            parts = [paragraph.text for paragraph in document.paragraphs]
            for table in document.tables:
                parts.extend(cell.text for row in table.rows for cell in row.cells)
            structured = {
                "document_type": "docx",
                "paragraph_count": len(document.paragraphs),
                "table_count": len(document.tables),
            }
        elif suffix in {"xlsx", "xlsm"}:
            from openpyxl import load_workbook  # type: ignore[import-untyped]

            workbook = load_workbook(path, read_only=True, data_only=True, keep_links=False)
            parts = []
            for sheet in workbook.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    parts.append("\t".join("" if value is None else str(value) for value in row))
            structured = {"document_type": suffix, "sheet_count": len(workbook.worksheets)}
            workbook.close()
        elif suffix == "pptx":
            from pptx import Presentation  # type: ignore[import-not-found]

            presentation = Presentation(str(path))
            parts = [
                shape.text
                for slide in presentation.slides
                for shape in slide.shapes
                if hasattr(shape, "text")
            ]
            structured = {"document_type": "pptx", "slide_count": len(presentation.slides)}
        elif suffix == "pdf":
            import fitz  # type: ignore[import-untyped, import-not-found]

            document = fitz.open(path)
            parts = [page.get_text("text") for page in document]
            structured = {"document_type": "pdf", "page_count": len(document)}
            document.close()
        else:
            return {"extraction_status": "UNSUPPORTED"}
        text = _bounded_text(parts, config)
        return {
            "normalized_text": text,
            "character_count": len(text),
            "word_count": len(text.split()),
            "normalized_text_hash": compute_normalized_text_hash(text),
            "structured_metadata": structured,
            "extraction_status": "OK",
        }
    except ImportError as exc:
        return {
            "extraction_status": "UNSUPPORTED",
            "extraction_error": f"optional parser unavailable: {exc.name}",
        }
    except Exception as exc:  # noqa: BLE001 - a malformed document is isolated and recorded, never fatal
        # A parser failure (corrupt zip container, bad XML, package error) is preserved as an
        # ERROR so it is never mistaken for clutter.  The failure is isolated to one file.
        return {"extraction_status": "ERROR", "extraction_error": f"{type(exc).__name__}: {exc}"}


def extract_docx_metadata(path: Path, config):
    return _structured_document(path, "docx", config)


def extract_xlsx_metadata(path: Path, config):
    return _structured_document(path, "xlsx", config)


def extract_pptx_metadata(path: Path, config):
    return _structured_document(path, "pptx", config)


def extract_pdf_metadata(path: Path, config):
    return _structured_document(path, "pdf", config)


def extract_document(path: Path, file_type: str, config):
    # ``file_type`` may arrive as ``.txt`` (from ``Path.suffix``) or ``txt``; normalize the
    # leading dot so plaintext extraction is reachable regardless of the caller's convention.
    normalized_type = file_type.lower().lstrip(".")
    return (
        extract_plaintext_metadata(path, config)
        if normalized_type in {"txt", "md", "csv", "rst", "log", "text"}
        else _structured_document(path, normalized_type, config)
    )


def run_document_analysis(database, config, scope=None, job_id=None):
    from .registry import run_content_analysis

    return run_content_analysis(database, config, "documents", job_id=job_id)
